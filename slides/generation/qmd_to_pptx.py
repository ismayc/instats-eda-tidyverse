import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

import re

def parse_qmd_slides(qmd_path):
    """Parses a Quarto/Markdown file and returns a list of slide dicts."""
    text = Path(qmd_path).read_text(encoding="utf-8")
    raw_slides = text.strip().split("\n---\n")
    parsed = []

    for block in raw_slides:
        lines = block.strip().split("\n")
        title, bullets, notes = "", [], []
        in_notes = False

        for line in lines:
            if line.startswith("## "):
                title = line.replace("## ", "").strip()
            elif line.startswith("::: {.notes}"):
                in_notes = True
            elif line.startswith(":::") and in_notes:
                in_notes = False
            elif in_notes:
                notes.append(line.strip())
            elif re.match(r"^\s*-\s", line):
                level = (len(line) - len(line.lstrip(" "))) // 2  # estimate indent level
                bullet_text = re.sub(r"^\s*-\s", "", line).strip()
                bullets.append((level, bullet_text))


        parsed.append({
            "title": title,
            "bullets": bullets,
            "notes": "\n".join(notes).strip()
        })

    return parsed

def format_text_with_backticks(text, paragraph, font_size, font_color):
    """
    Adds text to a paragraph with special formatting for backtick-wrapped inline code.
    Example: "This is `code`." → "This is" (normal), "`code`" (Inconsolata)
    """
    parts = re.split(r"(`[^`]+`)", text)
    for part in parts:
        run = paragraph.add_run()
        run.text = part.replace("`", "")
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.name = "Inconsolata" if part.startswith("`") and part.endswith("`") else None

def append_slides_to_template(template_path, slide_data, output_path):
    """Appends slides from parsed QMD to an existing PPTX template."""
    prs = Presentation(template_path)
    content_layout = prs.slide_layouts[1]  # Assumes 'Title and Content'

    bullet_color = RGBColor(0x4C, 0x91, 0xAF)
    title_color = RGBColor(0, 51, 102)

    for slide in slide_data:
        new_slide = prs.slides.add_slide(content_layout)

        # Title
        title_shape = new_slide.shapes.title
        title_shape.text_frame.clear()
        title_paragraph = title_shape.text_frame.paragraphs[0]
        format_text_with_backticks(slide["title"], title_paragraph, Pt(48), title_color)

        # Bullet content
        content_shape = next(
            (s for s in new_slide.shapes if s.has_text_frame and s != title_shape), None
        )
        if content_shape:
            tf = content_shape.text_frame
            tf.clear()
        for bullet in slide["bullets"]:
            if isinstance(bullet, tuple):
                level, text = bullet
            else:
                level, text = 0, bullet  # for backward compatibility
            para = tf.add_paragraph()
            para.level = level
            format_text_with_backticks(text, para, Pt(28), bullet_color)


        # Notes
        notes_frame = new_slide.notes_slide.notes_text_frame
        notes_frame.text = slide["notes"]

    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")

# --- USAGE EXAMPLE ---

if __name__ == "__main__":
    # Set paths (modify these as needed)
    qmd_file = "session6_slides_output.qmd"
    pptx_template = "template.pptx"
    output_file = "session6_case-study_slides.pptx"

    # Parse and generate
    slides = parse_qmd_slides(qmd_file)
    append_slides_to_template(pptx_template, slides, output_file)
