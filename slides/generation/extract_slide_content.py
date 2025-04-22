# pip install python-pptx
# reticulate::py_install("python-pptx")

# This script extracts the content from a PowerPoint presentation and formats it into Markdown.
# extract_slide_content.py
# Usage: call via reticulate with pptx_file and output_path already defined in R

from pptx import Presentation
import os

# These variables should be passed in from R before calling source_python()
# e.g., in R:
#   reticulate::py_run_string("pptx_file = '...'; output_path = '...'")
#   reticulate::source_python("extract_slide_content.py")

# Check if pptx_file and output_path exist
if not os.path.exists(pptx_file):
    raise FileNotFoundError(f"PowerPoint file not found: {pptx_file}")

# Create output directory if it doesn’t exist
os.makedirs(os.path.dirname(output_path), exist_ok=True)

# Load presentation
prs = Presentation(pptx_file)
markdown_output = []

# Process slides
for i, slide in enumerate(prs.slides, start=1):
    title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.has_text_frame else f"Slide {i}"

    # Extract bullet content
    bullets = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    bullets.append(f"- {text}")

    # Extract speaker notes
    notes_text = ""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            notes_text = notes_frame.text.strip()

    # Write slide to markdown
    markdown_output.append(f"## {title}\n")
    if bullets:
        markdown_output.extend(bullets)
    if notes_text:
        markdown_output.append(f"\n::: {{.notes}}\n{notes_text}\n:::\n")
    markdown_output.append("\n---\n")

# Write output
with open(output_path, "w", encoding="utf-8") as f:
    f.write("\n".join(markdown_output))
